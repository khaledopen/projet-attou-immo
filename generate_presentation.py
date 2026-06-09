import sys
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation and set 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette de couleurs - Thème moderne, épuré et premium
COLOR_BG = RGBColor(248, 250, 252)       # Slate 50 (Gris-bleu très clair)
COLOR_NAVY = RGBColor(15, 23, 42)       # Slate 900 (Bleu nuit moderne profond)
COLOR_BLUE = RGBColor(2, 132, 199)      # Sky 600 (Bleu actif vif)
COLOR_LIGHT_BLUE = RGBColor(224, 242, 254) # Sky 100 (Bleu ciel de fond doux)
COLOR_ORANGE = RGBColor(249, 115, 22)   # Orange 500 (Orange de marque chaleureux)
COLOR_WHITE = RGBColor(255, 255, 255)   # White
COLOR_GRAY = RGBColor(100, 116, 139)     # Slate 500 (Texte neutre subtil)
COLOR_BORDER = RGBColor(226, 232, 240)   # Slate 200 (Lignes de grille douces)

def set_slide_background(slide, color):
    """Sets a solid color background for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height, font_name="Calibri", font_size=16, font_color=COLOR_NAVY, bold=False, align=PP_ALIGN.LEFT):
    """Fonction d'aide pour ajouter un bloc de texte standardisé."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    return txBox

def add_paragraph(tf, text, font_size=14, font_color=COLOR_NAVY, bold=False, space_before=10):
    """Ajoute un paragraphe à un bloc de texte existant."""
    p = tf.add_paragraph()
    p.text = str(text)
    p.font.name = "Calibri"
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_slide_header(slide, title, category="ATTOUHOME - RAPPORT DE DÉVELOPPEMENT"):
    """Ajoute un en-tête premium standardisé aux diapositives normales."""
    # Étiquette de catégorie
    add_textbox(slide, category.upper(), Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3), font_size=10, font_color=COLOR_BLUE, bold=True)
    # Titre de la diapositive
    add_textbox(slide, title, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6), font_size=24, font_color=COLOR_NAVY, bold=True)

def add_card_shape(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_BORDER):
    """Crée un joli conteneur de carte pour structurer le contenu de la diapositive."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color if border_color else bg_color
    shape.line.width = Pt(1)
    return shape

# ==============================================================================
# DIAPOSITIVE 1 : DIAPOSITIVE DE COUVERTURE (Thème bleu nuit premium)
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, COLOR_NAVY)

# Grand élément de fond stylisé
accent_card = add_card_shape(slide1, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2), bg_color=COLOR_BLUE, border_color=None)

# Main Title
title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AttouHome"
p.font.name = "Calibri"
p.font.size = Pt(64)
p.font.color.rgb = COLOR_WHITE
p.font.bold = True

p_sub = tf.add_paragraph()
p_sub.text = "Révolutionner la Recherche & Gestion Immobilière"
p_sub.font.name = "Calibri"
p_sub.font.size = Pt(28)
p_sub.font.color.rgb = COLOR_BLUE
p_sub.font.bold = True
p_sub.space_before = Pt(15)

# Metadata Info
meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(11.0), Inches(1.5))
tf_meta = meta_box.text_frame
p_meta1 = tf_meta.paragraphs[0]
p_meta1.text = "Rapport Technique de Synthèse des Développements"
p_meta1.font.name = "Calibri"
p_meta1.font.size = Pt(14)
p_meta1.font.color.rgb = COLOR_GRAY
p_meta1.font.bold = True

p_meta2 = tf_meta.add_paragraph()
p_meta2.text = "Applications Mobiles (Tenant & Owner), Console d'Administration Web, Résilience Réseau & Synchronisations temps réel."
p_meta2.font.name = "Calibri"
p_meta2.font.size = Pt(12)
p_meta2.font.color.rgb = COLOR_GRAY
p_meta2.space_before = Pt(5)

# ==============================================================================
# DIAPOSITIVE 2 : VUE D'ENSEMBLE DE L'ARCHITECTURE 360° (Diapositive claire, mise en page multi-cartes)
# ==============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, COLOR_BG)
add_slide_header(slide2, "Une Architecture Multi-Plateforme Complète")

# Carte 1 : Backend & Prisma
add_card_shape(slide2, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide2, "⚡ CORE ENGINE & BASE", Inches(1.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=13, font_color=COLOR_BLUE, bold=True)
add_textbox(slide2, "Backend Node.js, Express & Prisma", Inches(1.1), Inches(2.3), Inches(3.0), Inches(0.6), font_size=18, font_color=COLOR_NAVY, bold=True)
desc1_box = slide2.shapes.add_textbox(Inches(1.1), Inches(3.0), Inches(3.0), Inches(3.0))
tf1 = desc1_box.text_frame
tf1.word_wrap = True
tf1.paragraphs[0].text = "• Architecture API REST sécurisée"
tf1.paragraphs[0].font.size = Pt(13)
tf1.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf1, "• Schémas de base de données PostgreSQL robustes managés via Prisma", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf1, "• Intégration en temps réel des services de modération d'annonces", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf1, "• Notifications Push & tunnels WebSocket actifs pour la messagerie instantanée", font_size=13, font_color=COLOR_GRAY)

# Carte 2 : Applications mobiles (Locataire & Propriétaire)
add_card_shape(slide2, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide2, "📱 COMPAGNON MOBILE", Inches(5.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=13, font_color=COLOR_ORANGE, bold=True)
add_textbox(slide2, "Applications iOS & Android", Inches(5.1), Inches(2.3), Inches(3.0), Inches(0.6), font_size=18, font_color=COLOR_NAVY, bold=True)
desc2_box = slide2.shapes.add_textbox(Inches(5.1), Inches(3.0), Inches(3.0), Inches(3.0))
tf2 = desc2_box.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = "• Apps compilées via Expo & React Native"
tf2.paragraphs[0].font.size = Pt(13)
tf2.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf2, "• Tenant App : Recherche premium, favoris locaux offline & prise de visites", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf2, "• Owner App : Dépôt d'annonces rapide avec sélection et upload d'images", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf2, "• Géolocalisation dynamique inverse pour le repérage auto des coordonnées", font_size=13, font_color=COLOR_GRAY)

# Carte 3 : Panneau web d'administration
add_card_shape(slide2, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide2, "🖥️ PILOTAGE CENTRAL", Inches(9.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=13, font_color=COLOR_BLUE, bold=True)
add_textbox(slide2, "IvoireAdmin Dashboard", Inches(9.1), Inches(2.3), Inches(3.0), Inches(0.6), font_size=18, font_color=COLOR_NAVY, bold=True)
desc3_box = slide2.shapes.add_textbox(Inches(9.1), Inches(3.0), Inches(3.0), Inches(3.0))
tf3 = desc3_box.text_frame
tf3.word_wrap = True
tf3.paragraphs[0].text = "• Panel Web d'administration moderne en React"
tf3.paragraphs[0].font.size = Pt(13)
tf3.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf3, "• Palette Navy/Orange corporate unifiée", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf3, "• Gestion des utilisateurs, modération des biens signalés & blocages instantanés", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf3, "• Statistiques dynamiques sur l'activité des visites et taux d'engagement", font_size=13, font_color=COLOR_GRAY)

# ==============================================================================
# DIAPOSITIVE 3 : RÉSILIENCE & AUTOMATISATION DU TUNNEL LOCAL (Diapositive claire, mise en page divisée en 2 colonnes)
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, COLOR_BG)
add_slide_header(slide3, "Zéro Coupure : Résilience Réseau & Tunneling Automatique")

# Colonne gauche (Concept & Tunnel)
add_card_shape(slide3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
add_textbox(slide3, "🔗 AUTOMATION DES ADRESSES IP", Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide3, "Mise à jour en temps réel des URLs", Inches(1.1), Inches(2.3), Inches(5.0), Inches(0.5), font_size=18, font_color=COLOR_NAVY, bold=True)
left_box = slide3.shapes.add_textbox(Inches(1.1), Inches(2.9), Inches(5.0), Inches(3.2))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.paragraphs[0].text = "Les tunnels Localtunnel/Serveo changent d'URL à chaque démarrage du serveur. Pour éviter les pannes ou les configurations manuelles laborieuses :"
tf_left.paragraphs[0].font.size = Pt(13)
tf_left.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_left, "✔ Capture automatique : Un script capture l'URL à la volée au démarrage.", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf_left, "✔ Synchronisation globale : Injection automatique de l'URL dans les fichiers .env de l'application Owner et Tenant simultanément.", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf_left, "✔ Zéro Friction : Les téléphones de test (iOS/Android) se connectent instantanément sans aucune configuration.", font_size=13, font_color=COLOR_GRAY)

# Colonne droite (Réseau & Tentatives)
add_card_shape(slide3, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
add_textbox(slide3, "🛡️ RESILIENCE AUX TENTATIVES RÉSEAU", Inches(7.1), Inches(1.9), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_ORANGE, bold=True)
add_textbox(slide3, "Intercepteurs Globaux Axios & Retries", Inches(7.1), Inches(2.3), Inches(5.0), Inches(0.5), font_size=18, font_color=COLOR_NAVY, bold=True)
right_box = slide3.shapes.add_textbox(Inches(7.1), Inches(2.9), Inches(5.0), Inches(3.2))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.paragraphs[0].text = "Les tunnels réseau de développement subissent souvent des pertes de paquets ou des lenteurs temporaires (erreurs 502/504) :"
tf_right.paragraphs[0].font.size = Pt(13)
tf_right.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_right, "✔ AxiosInstance standardisé : Les appels directs à 'axios' ont été remplacés par une instance réseau blindée commune aux applications.", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf_right, "✔ Gestionnaire de réessais (Retries) : Tentatives de reconnexion auto (jusqu'à 3 essais) avec un délai exponentiel avant de remonter une erreur.", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf_right, "✔ Upload d'images blindé : Publication stable des annonces même avec des fichiers volumineux.", font_size=13, font_color=COLOR_GRAY)

# ==============================================================================
# DIAPOSITIVE 4 : L'ÉCRAN D'EXPLORATION PREMIUM RECONÇU (Démonstration visuelle des fonctionnalités)
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, COLOR_BG)
add_slide_header(slide4, "L'Interface Explorer : Une Révolution Visuelle Moderne")

# Côté gauche : Liste des fonctionnalités
add_card_shape(slide4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
add_textbox(slide4, "🎨 EXPERIENCE PREMIUM WOW", Inches(1.1), Inches(1.8), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide4, "Fonctionnalités Clés Implémentées", Inches(1.1), Inches(2.2), Inches(5.0), Inches(0.5), font_size=18, font_color=COLOR_NAVY, bold=True)

feat_box = slide4.shapes.add_textbox(Inches(1.1), Inches(2.8), Inches(5.0), Inches(3.2))
tf_feat = feat_box.text_frame
tf_feat.word_wrap = True
tf_feat.paragraphs[0].text = "• En-tête Premium personnalisé : Salutations 'Bonjour 👋' et titre clair unifié à côté d'une cloche de notifications aux proportions parfaites."
tf_feat.paragraphs[0].font.size = Pt(12)
tf_feat.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_feat, "• Recherche & Filtres intégrés : Barre de recherche épurée couplée à un bouton de filtre de taille 44x44 élégant aux couleurs de la marque.", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_feat, "• Pills de sélection : Carrousel horizontal fluide (Tous, Appartement, Studio, Maison, Villa, Chambre) avec retour visuel actif bleu vif.", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_feat, "• Géolocalisation dynamique inversée : Localisation temps réel de l'utilisateur affichée dynamiquement (GPS vers Ville/Quartier).", font_size=12, font_color=COLOR_GRAY)

# Côté droit : Maquette visuelle de la carte de propriété premium
add_card_shape(slide4, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8), bg_color=COLOR_WHITE)
# En-tête de carte / Espace réservé pour l'image en bleu clair
add_card_shape(slide4, Inches(7.1), Inches(1.9), Inches(5.0), Inches(2.2), bg_color=COLOR_LIGHT_BLUE, border_color=None)
# Badge Coup de cœur dans la maquette
add_card_shape(slide4, Inches(7.3), Inches(2.1), Inches(1.8), Inches(0.35), bg_color=COLOR_BLUE, border_color=None)
add_textbox(slide4, "✦ Coup de cœur", Inches(7.3), Inches(2.12), Inches(1.8), Inches(0.35), font_size=9, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
# Badge cœur de favori dans la maquette
add_card_shape(slide4, Inches(11.3), Inches(2.1), Inches(0.5), Inches(0.5), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
add_textbox(slide4, "❤️", Inches(11.3), Inches(2.15), Inches(0.5), Inches(0.5), font_size=11, align=PP_ALIGN.CENTER)
# Photo counter
add_card_shape(slide4, Inches(11.1), Inches(3.6), Inches(0.8), Inches(0.3), bg_color=COLOR_NAVY, border_color=None)
add_textbox(slide4, "📷 3", Inches(11.1), Inches(3.6), Inches(0.8), Inches(0.3), font_size=8, font_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Mockup Details Content
add_textbox(slide4, "Appartement lumineux avec terrasse", Inches(7.1), Inches(4.3), Inches(3.6), Inches(0.3), font_size=13, font_color=COLOR_NAVY, bold=True)
add_textbox(slide4, "⭐ 4.8", Inches(11.0), Inches(4.3), Inches(1.1), Inches(0.3), font_size=12, font_color=COLOR_ORANGE, bold=True, align=PP_ALIGN.RIGHT)
add_textbox(slide4, "📍 Cocody · Abidjan", Inches(7.1), Inches(4.6), Inches(5.0), Inches(0.3), font_size=10, font_color=COLOR_GRAY)
# Ligne de séparation du pied de page
add_card_shape(slide4, Inches(7.1), Inches(5.0), Inches(5.0), Inches(0.02), bg_color=COLOR_BORDER, border_color=None)
# Specs
add_textbox(slide4, "🛌 2 ch.   📐 65 m²", Inches(7.1), Inches(5.1), Inches(2.5), Inches(0.3), font_size=11, font_color=COLOR_GRAY, bold=True)
add_textbox(slide4, "1 850 000 FCFA/mois", Inches(9.4), Inches(5.1), Inches(2.7), Inches(0.3), font_size=12, font_color=COLOR_BLUE, bold=True, align=PP_ALIGN.RIGHT)

# ==============================================================================
# DIAPOSITIVE 5 : PROTECTION CONTRE LE DÉBORDEMENT & STABILITÉ (Mise en page technique)
# ==============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, COLOR_BG)
add_slide_header(slide5, "Technique : Système Anti-Débordement de l'Interface")

# Les 3 piliers de la stabilité
add_card_shape(slide5, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide5, "1. FLEXWRAP DYNAMIQUE", Inches(1.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide5, "Adaptabilité aux Prix", Inches(1.1), Inches(2.3), Inches(3.0), Inches(0.5), font_size=17, font_color=COLOR_NAVY, bold=True)
tf_stab1 = slide5.shapes.add_textbox(Inches(1.1), Inches(2.9), Inches(3.0), Inches(3.0)).text_frame
tf_stab1.word_wrap = True
tf_stab1.paragraphs[0].text = "Les montants en FCFA atteignent régulièrement des millions, exigeant un espace horizontal important :"
tf_stab1.paragraphs[0].font.size = Pt(12)
tf_stab1.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_stab1, "✔ Le pied de page de la carte autorise dorénavant un enveloppement automatique ('flexWrap: wrap').", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_stab1, "✔ Si le prix et les caractéristiques entrent en collision, le prix descend sur une seconde ligne en toute sécurité.", font_size=12, font_color=COLOR_GRAY)

add_card_shape(slide5, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide5, "2. ALIGNEMENT AUTOMATIQUE", Inches(5.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=12, font_color=COLOR_ORANGE, bold=True)
add_textbox(slide5, "Astuces Flexbox Modernes", Inches(5.1), Inches(2.3), Inches(3.0), Inches(0.5), font_size=17, font_color=COLOR_NAVY, bold=True)
tf_stab2 = slide5.shapes.add_textbox(Inches(5.1), Inches(2.9), Inches(3.0), Inches(3.0)).text_frame
tf_stab2.word_wrap = True
tf_stab2.paragraphs[0].text = "Garantir un alignement visuel parfait sur les grilles, même en cas de retour à la ligne :"
tf_stab2.paragraphs[0].font.size = Pt(12)
tf_stab2.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_stab2, "✔ Utilisation de 'marginLeft: auto' sur le conteneur de prix.", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_stab2, "✔ Si le prix reste sur la même ligne, il s'aligne à droite. S'il descend à la ligne, il se pousse automatiquement vers la droite, conservant un aspect de tableau extrêmement chic.", font_size=12, font_color=COLOR_GRAY)

add_card_shape(slide5, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.8))
add_textbox(slide5, "3. LIMITATIONS PHYSIQUES", Inches(9.1), Inches(1.9), Inches(3.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide5, "Ellipses & Réduction", Inches(9.1), Inches(2.3), Inches(3.0), Inches(0.5), font_size=17, font_color=COLOR_NAVY, bold=True)
tf_stab3 = slide5.shapes.add_textbox(Inches(9.1), Inches(2.9), Inches(3.0), Inches(3.0)).text_frame
tf_stab3.word_wrap = True
tf_stab3.paragraphs[0].text = "Gestion rigoureuse des longueurs de textes imprévues saisies par les propriétaires :"
tf_stab3.paragraphs[0].font.size = Pt(12)
tf_stab3.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_stab3, "✔ Blocage du titre et de la localisation à 1 seule ligne maximum avec troncation ('numberOfLines={1}').", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_stab3, "✔ Saisie modifiée pour forcer la séparation des blocs textuels au lieu des composants textes imbriqués bogués sous Android.", font_size=12, font_color=COLOR_GRAY)

# ==============================================================================
# DIAPOSITIVE 6 : AUTHENTICATION, VISITES & MODÉRATION (Diapositive de résumé)
# ==============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, COLOR_BG)
add_slide_header(slide6, "Sécurité, Prises de Rendez-vous & Modération des Biens")

# Carte gauche : Sécurité & Authentification
add_card_shape(slide6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
add_textbox(slide6, "🔒 PROTECTION DES DONNÉES", Inches(1.1), Inches(1.9), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide6, "Accès Limité & Comptes Suspendus", Inches(1.1), Inches(2.3), Inches(5.0), Inches(0.5), font_size=18, font_color=COLOR_NAVY, bold=True)
tf_sec = slide6.shapes.add_textbox(Inches(1.1), Inches(2.9), Inches(5.0), Inches(3.2)).text_frame
tf_sec.word_wrap = True
tf_sec.paragraphs[0].text = "• Parcours de Conversion : Les détails complets des biens, équipements et favoris sont ouverts sans connexion pour maximiser l'intérêt. L'authentification n'est exigée qu'à la demande de visite ou au contact."
tf_sec.paragraphs[0].font.size = Pt(12)
tf_sec.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_sec, "• Modération active : Les locataires peuvent signaler une annonce douteuse en saisissant un motif détaillé via un KeyboardAvoidingView protecteur.", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_sec, "• Blocages instantanés : L'extension Prisma permet de marquer un compte comme SUSPENDED dans le panel d'administration. La connexion de ce dernier et ses requêtes sont immédiatement invalidées côté backend.", font_size=12, font_color=COLOR_GRAY)

# Carte droite : Visites en temps réel & Chat
add_card_shape(slide6, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
add_textbox(slide6, "📅 INTERACTION EN DIRECT", Inches(7.1), Inches(1.9), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_ORANGE, bold=True)
add_textbox(slide6, "Messagerie & Calendrier de Visites", Inches(7.1), Inches(2.3), Inches(5.0), Inches(0.5), font_size=18, font_color=COLOR_NAVY, bold=True)
tf_vis = slide6.shapes.add_textbox(Inches(7.1), Inches(2.9), Inches(5.0), Inches(3.2)).text_frame
tf_vis.word_wrap = True
tf_vis.paragraphs[0].text = "• Planification interactive : Un calendrier intégré permet au locataire de sélectionner un jour libre pour visiter la propriété. La demande est immédiatement notifiée à l'application Owner."
tf_vis.paragraphs[0].font.size = Pt(12)
tf_vis.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_vis, "• Système anti-doublon : Les demandes de visites redondantes ou multiples sur la même annonce par le même compte sont bloquées en amont.", font_size=12, font_color=COLOR_GRAY)
add_paragraph(tf_vis, "• Messagerie Socket : Échange de messages cryptés et fluides entre le propriétaire et le locataire avec affichage des horaires et badges de notification non-lus sur l'onglet.", font_size=12, font_color=COLOR_GRAY)

# ==============================================================================
# DIAPOSITIVE 7 : RÉSUMÉ & PROCHAINES ÉTAPES (Thème séparateur bleu nuit)
# ==============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, COLOR_NAVY)

# Grande ligne d'accentuation verticale
add_card_shape(slide7, Inches(0.8), Inches(1.8), Inches(0.15), Inches(4.0), bg_color=COLOR_ORANGE, border_color=None)

# Zone de texte du côté gauche
add_textbox(slide7, "CONCLUSION & FUTURS HORIZONS", Inches(1.2), Inches(1.7), Inches(5.0), Inches(0.4), font_size=12, font_color=COLOR_BLUE, bold=True)
add_textbox(slide7, "Bilan des Avancées", Inches(1.2), Inches(2.1), Inches(5.0), Inches(0.6), font_size=28, font_color=COLOR_WHITE, bold=True)

summary_box = slide7.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(5.0), Inches(3.5))
tf_sum = summary_box.text_frame
tf_sum.word_wrap = True
tf_sum.paragraphs[0].text = "L'application AttouHome dispose d'une infrastructure technique solide, hautement résiliente aux contraintes du réseau de développement local. Les interfaces de l'application mobile et de la console d'administration sont unifiées, stables, dynamiques et conçues selon les meilleurs standards d'ergonomie (UI/UX)."
tf_sum.paragraphs[0].font.size = Pt(13)
tf_sum.paragraphs[0].font.color.rgb = COLOR_GRAY

# Zone de texte côté droit - Prochaines étapes
add_textbox(slide7, "PROCHAINES ÉTAPES DE PRODUCTION", Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=12, font_color=COLOR_ORANGE, bold=True)
add_textbox(slide7, "Chantiers Prioritaires", Inches(6.8), Inches(2.1), Inches(5.5), Inches(0.6), font_size=28, font_color=COLOR_WHITE, bold=True)

steps_box = slide7.shapes.add_textbox(Inches(6.8), Inches(2.9), Inches(5.5), Inches(3.5))
tf_steps = steps_box.text_frame
tf_steps.word_wrap = True
tf_steps.paragraphs[0].text = "1. Phase de Recette (UAT) : Simulation de visites à grande échelle avec des comptes fictifs."
tf_steps.paragraphs[0].font.size = Pt(13)
tf_steps.paragraphs[0].font.color.rgb = COLOR_GRAY
add_paragraph(tf_steps, "2. Déploiement Cloud Production : Migration de la base PostgreSQL locale et du backend vers AWS/Heroku pour s'affranchir définitivement des tunnels de dev.", font_size=13, font_color=COLOR_GRAY)
add_paragraph(tf_steps, "3. Build natif de test : Exportation des fichiers IPA (iOS) et APK (Android) via EAS Build d'Expo pour distribution aux bêta-testeurs.", font_size=13, font_color=COLOR_GRAY)

# Sauvegarder la présentation
prs.save("AttouHome_Rapport_Developpement.pptx")
print("Presentation successfully created as AttouHome_Rapport_Developpement.pptx!")
